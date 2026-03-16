import os
import fitz  # PyMuPDF
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import cv2
from pptx import Presentation

from sentence_transformers import SentenceTransformer, util

import nltk
nltk.download('punkt', quiet=True)

# Use path based on launching from project root!
DATA_DIR = r"C:\Users\srima\Desktop\A3_SETS\occupational_safety_\data\sample\샘플"
OUT_DIR = "output/parsed/"
LOG_FILE = "output/parse_log.txt"
os.makedirs(OUT_DIR, exist_ok=True)

def log(message):
    with open(LOG_FILE, "a", encoding="utf-8") as logf:
        logf.write(message + "\n")
    print(message)

def parse_pdf(filepath):
    text = ""
    pdf = fitz.open(filepath)
    for page in pdf:
        text += page.get_text()
    if not text.strip():
        log(f"[PDF][OCR] Scanned: {filepath}")
        images = convert_from_path(filepath)
        for i, img in enumerate(images):
            text += f"\n--- Page {i+1} ---\n"
            text += pytesseract.image_to_string(img, lang="kor+eng")
    else:
        log(f"[PDF][Digital] Parsed: {filepath}")
    return text

def parse_txt(filepath):
    try:
        with open(filepath, encoding="cp949") as f:
            text = f.read()
    except UnicodeDecodeError:
        with open(filepath, encoding="utf-8") as f:
            text = f.read()
    log(f"[TXT] Parsed: {filepath}")
    return text

def parse_img(filepath):
    img = Image.open(filepath)
    text = pytesseract.image_to_string(img, lang="kor+eng")
    log(f"[IMG][OCR] Parsed: {filepath}")
    return text

def parse_video(filepath):
    text = ""
    cap = cv2.VideoCapture(filepath)
    fps = int(cap.get(cv2.CAP_PROP_FPS)) or 1
    count = 0
    frame_num = 0
    while True:
        ret, frame = cap.read()
        if not ret:
            break
        if frame_num % max(1, fps) == 0:
            img_pil = Image.fromarray(frame)
            text += f"\n--- Frame {frame_num} ---\n"
            text += pytesseract.image_to_string(img_pil, lang="kor+eng")
            count += 1
        frame_num += 1
    cap.release()
    log(f"[VIDEO][OCR] Parsed: {filepath} ({count} frames OCRed)")
    return text

def parse_pptx(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n\n"
    log(f"[PPTX] Parsed: {filepath}")
    return text

def chunk_text(text):
    # Classic paragraph-based chunking
    return [chunk.strip() for chunk in text.split('\n\n') if chunk.strip()]

def semantic_chunk_text(text, threshold=0.75):
    sentences = nltk.sent_tokenize(text)
    if not sentences:
        return []
    model = SentenceTransformer('paraphrase-multilingual-MiniLM-L12-v2')
    embeddings = model.encode(sentences)
    chunks = []
    current_chunk = [sentences[0]]

    for i in range(1, len(sentences)):
        sim = util.cos_sim(embeddings[i - 1], embeddings[i]).item()
        if sim < threshold:
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentences[i]]
        else:
            current_chunk.append(sentences[i])
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    return chunks

if __name__ == "__main__":
    log("==== Batch Parsing Started ====")
    files = os.listdir(DATA_DIR)
    for fname in files:
        fpath = os.path.join(DATA_DIR, fname)
        ext = fname.lower().split('.')[-1]
        if ext == "pdf":
            text = parse_pdf(fpath)
        elif ext == "txt":
            text = parse_txt(fpath)
        elif ext in {"jpg", "jpeg", "png", "bmp"}:
            text = parse_img(fpath)
        elif ext in {"mp4", "avi", "mov"}:
            text = parse_video(fpath)
        elif ext == "pptx":
            text = parse_pptx(fpath)
        else:
            log(f"[SKIP] Unsupported file: {fpath}")
            continue

        if not text.strip():
            log(f"[EMPTY] No text extracted: {fpath}")
            continue

        # Save classic (paragraph) chunks
        chunks = chunk_text(text)
        out_file = os.path.join(OUT_DIR, fname + ".chunks.txt")
        with open(out_file, "w", encoding="utf-8") as f:
            for i, chunk in enumerate(chunks):
                f.write(f"[Chunk {i + 1}]\n{chunk}\n\n")
        log(f"[DONE] Saved {out_file}; {len(chunks)} chunks")

        # Save semantic chunks
        sem_chunks = semantic_chunk_text(text)
        sem_out_file = os.path.join(OUT_DIR, fname + ".semantic_chunks.txt")
        with open(sem_out_file, "w", encoding="utf-8") as f:
            for i, chunk in enumerate(sem_chunks):
                f.write(f"[Semantic Chunk {i + 1}]\n{chunk}\n\n")
        log(f"[DONE][Semantic] Saved {sem_out_file}; {len(sem_chunks)} semantic chunks")

    log("==== Batch Parsing Finished ====")
